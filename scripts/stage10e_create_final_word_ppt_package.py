from pathlib import Path
import pandas as pd
import numpy as np

from docx import Document
from docx.shared import Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches, Pt

REPO = Path.home() / "ISPY2-3DGCNN"
REPORTS = REPO / "reports"
TABLES = REPORTS / "tables"
FIGS = REPORTS / "figures"
FINAL = REPORTS / "final_package"
VIEWS = REPORTS / "table_views"

for p in [REPORTS, TABLES, FIGS, FINAL, VIEWS]:
    p.mkdir(parents=True, exist_ok=True)

def read_csv(path):
    path = Path(path)
    if path.exists():
        return pd.read_csv(path, dtype=str).fillna("")
    return pd.DataFrame()

def safe_float(x):
    try:
        return float(x)
    except Exception:
        return np.nan

def best_from_table(path):
    df = read_csv(path)
    if df.empty:
        return {}
    d = df.copy()
    for c in ["AUROC", "AUPRC", "Brier_score"]:
        if c in d.columns:
            d[c] = pd.to_numeric(d[c], errors="coerce")
    d = d.sort_values(["AUROC", "AUPRC", "Brier_score"], ascending=[False, False, True])
    return d.iloc[0].to_dict()

def add_doc_heading(doc, text, level=1):
    doc.add_heading(text, level=level)

def add_doc_para(doc, text, bold_label=None):
    if bold_label:
        p = doc.add_paragraph()
        r = p.add_run(bold_label)
        r.bold = True
        p.add_run(" " + text)
    else:
        doc.add_paragraph(text)

def add_doc_image(doc, img_path, caption, width=6.0):
    img_path = Path(img_path)
    if img_path.exists():
        doc.add_picture(str(img_path), width=DocxInches(width))
        p = doc.add_paragraph(caption)
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER

def add_doc_table_from_df(doc, df, title, max_rows=8, max_cols=5):
    if df.empty:
        return
    doc.add_paragraph(title).runs[0].bold = True
    small = df.head(max_rows).iloc[:, :max_cols]
    table = doc.add_table(rows=1, cols=len(small.columns))
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, c in enumerate(small.columns):
        hdr[i].text = str(c)
    for _, row in small.iterrows():
        cells = table.add_row().cells
        for i, val in enumerate(row.tolist()):
            cells[i].text = str(val)
    doc.add_paragraph("")

def add_slide_title(slide, title, subtitle=None):
    shapes = slide.shapes
    box = shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.6))
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    if subtitle:
        box2 = shapes.add_textbox(Inches(0.45), Inches(0.9), Inches(12.0), Inches(0.4))
        tf2 = box2.text_frame
        tf2.text = subtitle
        tf2.paragraphs[0].font.size = Pt(14)

def add_slide_bullets(slide, title, bullets):
    add_slide_title(slide, title)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12.0), Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = str(b)
        p.font.size = Pt(18)
        p.level = 0

def add_slide_image(slide, title, img_path, caption=None):
    add_slide_title(slide, title)
    img_path = Path(img_path)
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.7), Inches(1.15), width=Inches(11.8))
    else:
        box = slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(10.5), Inches(1.0))
        box.text_frame.text = "Missing figure: " + img_path.name
    if caption:
        cap = slide.shapes.add_textbox(Inches(0.8), Inches(6.75), Inches(11.5), Inches(0.35))
        cap.text_frame.text = caption
        cap.text_frame.paragraphs[0].font.size = Pt(10)

# ---------------------------------------------------------
# Load final project information
# ---------------------------------------------------------
stage10d_story = read_csv(TABLES / "stage10d_final_thesis_story_summary.csv")
stage10d_caption = read_csv(TABLES / "stage10d_thesis_figure_caption_table.csv")
project_status = read_csv(TABLES / "phase8e_final_project_status.csv")
claims = read_csv(TABLES / "phase8e_safe_claims_for_paper.csv")
p9c_plan = read_csv(TABLES / "phase9c_next_research_plan.csv")

phase8c_best = best_from_table(TABLES / "phase8c_ftv_model_performance.csv")
phase5a_best = best_from_table(TABLES / "phase5a_model_calibration_metrics.csv")
phase9b_best = best_from_table(TABLES / "phase9b_temporal_delta_model_performance.csv")

dataset_summary = read_csv(TABLES / "stage10a_dataset_structure_summary.csv")
dataset_map = dict(zip(dataset_summary.get("metric", []), dataset_summary.get("value", [])))

dataset_size = dataset_map.get("dataset_size_du", "")
patients = dataset_map.get("patients", "")
series = dataset_map.get("series", "")
dicom_files = dataset_map.get("dicom_files_counted", "")

ftv_model = phase8c_best.get("model", "SVM_RBF")
ftv_auroc = safe_float(phase8c_best.get("AUROC", "nan"))
ftv_auprc = safe_float(phase8c_best.get("AUPRC", "nan"))
ftv_brier = safe_float(phase8c_best.get("Brier_score", "nan"))
ftv_n = phase8c_best.get("n_patients", "13")

base_model = phase5a_best.get("model", "")
base_auroc = safe_float(phase5a_best.get("AUROC", "nan"))
base_auprc = safe_float(phase5a_best.get("AUPRC", "nan"))
base_brier = safe_float(phase5a_best.get("Brier_score", "nan"))

temp_model = phase9b_best.get("model", "")
temp_auroc = safe_float(phase9b_best.get("AUROC", "nan"))
temp_auprc = safe_float(phase9b_best.get("AUPRC", "nan"))
temp_brier = safe_float(phase9b_best.get("Brier_score", "nan"))

# Figures
fig_master = FIGS / "340_stage10d_master_methods_panel.png"
fig_contact = FIGS / "341_stage10d_selected_case_contact_sheet.png"
fig_final = FIGS / "342_stage10d_final_pipeline_and_model_result.png"
fig_pipeline = FIGS / "270_phase8e_project_pipeline.png"
fig_metrics = FIGS / "260_phase8d_baseline_vs_ftv_metrics.png"
fig_graph = FIGS / "283_phase9a_graph_schema.png"
fig_decision = FIGS / "300_phase9c_model_decision.png"

# ---------------------------------------------------------
# Final Markdown report
# ---------------------------------------------------------
md = []
md.append("# Final Report Package")
md.append("")
md.append("## Title")
md.append("Pilot FTV-Based Prediction of Pathologic Complete Response in I-SPY2 Breast MRI with Tumor-Source Validation and a Future Temporal Graph Learning Framework")
md.append("")
md.append("## Main idea")
md.append("This project built a reproducible Cradle-to-GitHub pipeline for I-SPY2 breast MRI, tested DICOM SEG masks, found that the raw SEG path was not reliable enough for the main model, and then moved to official MRI/FTV support-data features for pCR prediction.")
md.append("")
md.append("## Main result")
md.append(f"The best current pilot model is {ftv_model} using official MRI/FTV support features. It achieved AUROC {ftv_auroc:.4f}, AUPRC {ftv_auprc:.4f}, and Brier score {ftv_brier:.4f} on {ftv_n} labeled patients.")
md.append("")
md.append("## Model decision")
md.append(f"The temporal-delta model used {temp_model} and reached AUROC {temp_auroc:.4f}. Since it was weaker than the FTV model, temporal GNN should be future work, not the main current model.")
md.append("")
md.append("## Safe claim")
md.append("Official MRI/FTV support-data features produced a stronger pilot pCR prediction signal than the earlier proxy and temporal-delta approaches.")
md.append("")
md.append("## Limitation")
md.append("This is pilot work only. It is not clinical validation because the strongest FTV model used a small labeled cohort.")
md.append("")
md.append("## Final outputs")
md.append("- final_package/ISPY2_FTV_pCR_Final_Report.docx")
md.append("- final_package/ISPY2_FTV_pCR_Final_Presentation.pptx")
md.append("- final_package/ISPY2_FTV_pCR_One_Page_Summary.md")
md.append("- reports/Stage10E_Final_Word_PowerPoint_Package_Report.md")
md.append("")
(FINAL / "ISPY2_FTV_pCR_One_Page_Summary.md").write_text("\n".join(md))

# ---------------------------------------------------------
# Create Word report
# ---------------------------------------------------------
doc = Document()
doc.add_heading("Pilot FTV-Based Prediction of Pathologic Complete Response in I-SPY2 Breast MRI", 0)
doc.add_paragraph("Tumor-source validation with a future temporal graph learning framework")

add_doc_heading(doc, "Abstract", 1)
add_doc_para(
    doc,
    "This pilot study built a reproducible Cradle-to-GitHub pipeline for I-SPY2 breast MRI. "
    "The project tested raw DICOM SEG masks, found that they were not reliable tumor-only masks for the main analysis, "
    "moved to official MRI/FTV support-data features, mapped pCR labels, and trained pilot pCR prediction models. "
    f"The strongest current result is an {ftv_model} model using FTV support features with AUROC {ftv_auroc:.4f}, "
    f"AUPRC {ftv_auprc:.4f}, and Brier score {ftv_brier:.4f}."
)

add_doc_heading(doc, "1. Research Question", 1)
add_doc_para(
    doc,
    "Can official MRI/FTV support-data features provide a stronger and more reliable pilot pCR prediction signal "
    "than features derived from weak DICOM SEG or proxy biomarker paths?"
)

add_doc_heading(doc, "2. Dataset and Computing Environment", 1)
add_doc_para(doc, f"Dataset size on Cradle: {dataset_size}.")
add_doc_para(doc, f"Patients in inventory: {patients}.")
add_doc_para(doc, f"Series counted: {series}.")
add_doc_para(doc, f"DICOM files counted: {dicom_files}.")

add_doc_heading(doc, "3. Tumor-Source Validation", 1)
add_doc_para(
    doc,
    "The DICOM SEG path was tested first. It became a diagnostic part of the study rather than the final modeling path. "
    "Many SEG files behaved like full-slice analysis masks. Fractional mask recovery was possible for a limited subset, "
    "but it did not provide enough labeled patients for reliable supervised modeling."
)

add_doc_heading(doc, "4. Main Modeling Path", 1)
add_doc_para(
    doc,
    "The official MRI/FTV support-data path was selected as the main current result because it gave the strongest pilot pCR prediction signal."
)
add_doc_table_from_df(doc, read_csv(TABLES / "phase8d_baseline_vs_ftv_model_comparison.csv"), "Table 1. Baseline vs FTV model comparison", max_rows=4, max_cols=8)

add_doc_heading(doc, "5. Main Result", 1)
add_doc_para(doc, f"Best FTV model: {ftv_model}.")
add_doc_para(doc, f"AUROC: {ftv_auroc:.4f}.")
add_doc_para(doc, f"AUPRC: {ftv_auprc:.4f}.")
add_doc_para(doc, f"Brier score: {ftv_brier:.4f}.")
add_doc_image(doc, fig_metrics, "Figure 1. Baseline vs FTV support-data model metrics.", width=6.0)

add_doc_heading(doc, "6. Visual Methods Output", 1)
add_doc_para(
    doc,
    "Stage 10D packaged the visual preprocessing, recovered tumor mask overlay, breast-pair formation, and contrast-transport proxy into thesis-ready figures."
)
add_doc_image(doc, fig_master, "Figure 2. Master methods panel.", width=6.0)
add_doc_image(doc, fig_final, "Figure 3. Final pipeline and model result figure.", width=6.0)

add_doc_heading(doc, "7. Temporal Graph Decision", 1)
add_doc_para(
    doc,
    "A graph-ready longitudinal FTV dataset was created. However, the temporal-delta model did not beat the FTV support-data model. "
    "Therefore, temporal GNN is future work, not the main current model."
)
add_doc_table_from_df(doc, read_csv(TABLES / "phase9c_model_decision_summary.csv"), "Table 2. Phase 9C model decision summary", max_rows=4, max_cols=7)
add_doc_image(doc, fig_graph, "Figure 4. Future graph-ready longitudinal representation.", width=6.0)

add_doc_heading(doc, "8. Limitations", 1)
add_doc_para(doc, "This is pilot work. The strongest model used a small labeled cohort and should not be described as clinical validation.")
add_doc_para(doc, "The obstacle proxy is a conceptual contrast-transport visualization. It is not a direct measurement of blood flow, milk flow, or fluid flow.")
add_doc_para(doc, "The temporal GNN direction needs more patients and more stable longitudinal tumor-region features.")

add_doc_heading(doc, "9. Future Work", 1)
add_doc_para(
    doc,
    "Future work should use a larger longitudinal FTV/PE/SER cohort and build a temporal graph model with stable patient-visit or tumor-region nodes. "
    "A future GNN should be compared against the current FTV SVM baseline before making stronger claims."
)

add_doc_heading(doc, "10. Recommended Thesis Figures", 1)
add_doc_table_from_df(doc, stage10d_caption, "Table 3. Figure caption map", max_rows=8, max_cols=3)

docx_path = FINAL / "ISPY2_FTV_pCR_Final_Report.docx"
doc.save(docx_path)

# ---------------------------------------------------------
# Create PowerPoint
# ---------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_title(
    slide,
    "Pilot FTV-Based pCR Prediction in I-SPY2 Breast MRI",
    "Tumor-source validation with a future temporal graph learning framework"
)

# Why this study
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bullets(slide, "Research problem", [
    "Breast MRI can monitor response during neoadjuvant therapy.",
    "Machine learning depends on reliable tumor features.",
    "Raw DICOM SEG masks were not reliable tumor-only masks in this pipeline.",
    "Official MRI/FTV support features became the stronger modeling path."
])

# Dataset
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bullets(slide, "Dataset and pipeline", [
    f"Dataset size on Cradle: {dataset_size}",
    f"Patients in inventory: {patients}",
    f"Series: {series}",
    f"DICOM files counted: {dicom_files}",
    "Workflow: inventory, SEG diagnosis, support-data FTV modeling, graph-ready dataset."
])

# Pipeline figure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_image(slide, "Project pipeline", fig_pipeline, "Overall project workflow.")

# Visual methods
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_image(slide, "Visual preprocessing and breast-pair methods", fig_master, "Stage 10D master methods panel.")

# Main result
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bullets(slide, "Main pilot result", [
    f"Best current model: {ftv_model}",
    f"FTV AUROC: {ftv_auroc:.4f}",
    f"FTV AUPRC: {ftv_auprc:.4f}",
    f"FTV Brier score: {ftv_brier:.4f}",
    f"Earlier baseline AUROC: {base_auroc:.4f}"
])

# Metrics figure
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_image(slide, "FTV model vs baseline", fig_metrics, "FTV support-data model gave the strongest pilot signal.")

# Temporal graph decision
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_image(slide, "Temporal graph decision", fig_decision, "Temporal graph learning is future work, not the main current model.")

# Future graph
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_image(slide, "Future temporal graph framework", fig_graph, "Patient visits become graph nodes; temporal edges track treatment change.")

# Limitations
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bullets(slide, "Limitations", [
    "Pilot cohort is small.",
    "The best FTV model used 13 labeled patients.",
    "This is not clinical validation.",
    "DICOM SEG masks were diagnostic, not the final tumor-source path.",
    "Temporal GNN needs more patients and stable longitudinal node identity."
])

# Final conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bullets(slide, "Conclusion", [
    "The project built a reproducible ISPY2 Cradle-to-GitHub pipeline.",
    "Tumor-source validation changed the research direction.",
    "Official MRI/FTV support features gave the strongest pilot pCR signal.",
    "Temporal graph learning remains the next research phase."
])

pptx_path = FINAL / "ISPY2_FTV_pCR_Final_Presentation.pptx"
prs.save(pptx_path)

# ---------------------------------------------------------
# Manifest and package report
# ---------------------------------------------------------
manifest = pd.DataFrame([
    {"file": str(docx_path.relative_to(REPO)), "type": "Word report", "purpose": "Final thesis or proposal report draft"},
    {"file": str(pptx_path.relative_to(REPO)), "type": "PowerPoint", "purpose": "Professor meeting or defense slides"},
    {"file": "reports/final_package/ISPY2_FTV_pCR_One_Page_Summary.md", "type": "Markdown summary", "purpose": "Short summary"},
])

manifest_path = TABLES / "stage10e_final_document_package_manifest.csv"
manifest.to_csv(manifest_path, index=False)

report = []
report.append("# Stage 10E Final Word and PowerPoint Package Report")
report.append("")
report.append("This stage creates the final Word report and PowerPoint presentation from the completed pipeline.")
report.append("")
report.append("## Main outputs")
report.append("")
report.append("- reports/final_package/ISPY2_FTV_pCR_Final_Report.docx")
report.append("- reports/final_package/ISPY2_FTV_pCR_Final_Presentation.pptx")
report.append("- reports/final_package/ISPY2_FTV_pCR_One_Page_Summary.md")
report.append("- reports/tables/stage10e_final_document_package_manifest.csv")
report.append("")
report.append("## Main story")
report.append("")
report.append("The main current result is the Phase 8C FTV support-data SVM_RBF model.")
report.append("")
report.append(f"Best FTV AUROC: {ftv_auroc:.4f}")
report.append(f"Best FTV AUPRC: {ftv_auprc:.4f}")
report.append(f"Best FTV Brier score: {ftv_brier:.4f}")
report.append("")
report.append("## Safe interpretation")
report.append("")
report.append("This is a pilot feasibility result, not clinical validation.")
report.append("")
report.append("## Next step")
report.append("")
report.append("Open the Word report and PowerPoint from reports/final_package. Edit wording with your professor's feedback.")

stage10e_report = REPORTS / "Stage10E_Final_Word_PowerPoint_Package_Report.md"
stage10e_report.write_text("\n".join(report))

# Dashboard/README update
dash = REPORTS / "Dashboard.md"
old_dash = dash.read_text() if dash.exists() else "# ISPY2 4D Atlas Dashboard\n"
dash_add = (
    "\n\n## Stage 10E final report and presentation\n\n"
    "- [Final Word/PPT package report](Stage10E_Final_Word_PowerPoint_Package_Report.md)\n"
    "- reports/final_package/ISPY2_FTV_pCR_Final_Report.docx\n"
    "- reports/final_package/ISPY2_FTV_pCR_Final_Presentation.pptx\n"
)
if "Stage 10E final report and presentation" not in old_dash:
    dash.write_text(old_dash.rstrip() + dash_add)

readme = REPO / "README.md"
old_readme = readme.read_text() if readme.exists() else "# ISPY2-3DGCNN\n"
readme_add = (
    "\n\n### Stage 10E: final Word report and PowerPoint\n\n"
    "Main outputs:\n\n"
    "- reports/final_package/ISPY2_FTV_pCR_Final_Report.docx\n"
    "- reports/final_package/ISPY2_FTV_pCR_Final_Presentation.pptx\n"
    "- reports/Stage10E_Final_Word_PowerPoint_Package_Report.md\n"
)
if "Stage 10E: final Word report and PowerPoint" not in old_readme:
    readme.write_text(old_readme.rstrip() + readme_add)

print("Stage 10E complete.")
print("Word report:", docx_path)
print("PowerPoint:", pptx_path)
print("Summary:", FINAL / "ISPY2_FTV_pCR_One_Page_Summary.md")
